#!/usr/bin/env python3
"""
ZipMart Investor Pitch Presentation Generator
Creates a professional PowerPoint presentation for investor pitches
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor

# Color scheme - Indian quick commerce theme (Orange & White)
PRIMARY_COLOR = RGBColor(255, 108, 0)      # Bright Orange
SECONDARY_COLOR = RGBColor(245, 245, 245)  # Off-white
TEXT_COLOR = RGBColor(40, 40, 40)          # Dark gray
WHITE = RGBColor(255, 255, 255)
ACCENT_GREEN = RGBColor(34, 177, 76)       # Green for positives
ACCENT_RED = RGBColor(220, 53, 69)         # Red for contrasts

def add_title_slide(prs, title, subtitle):
    """Add a title slide"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = PRIMARY_COLOR
    
    # Add title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(9), Inches(1.5))
    title_frame = title_box.text_frame
    title_frame.word_wrap = True
    title_p = title_frame.paragraphs[0]
    title_p.text = title
    title_p.font.size = Pt(66)
    title_p.font.bold = True
    title_p.font.color.rgb = WHITE
    
    # Add subtitle
    subtitle_box = slide.shapes.add_textbox(Inches(0.5), Inches(4.2), Inches(9), Inches(2))
    subtitle_frame = subtitle_box.text_frame
    subtitle_frame.word_wrap = True
    for line in subtitle.split('\n'):
        if subtitle_frame.paragraphs[0].text == '':
            p = subtitle_frame.paragraphs[0]
        else:
            p = subtitle_frame.add_paragraph()
        p.text = line
        p.font.size = Pt(28)
        p.font.color.rgb = WHITE
        p.space_before = Pt(6)

def add_content_slide(prs, title, content_type="bullets"):
    """Add a content slide with title"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
    
    # Add title bar
    title_shape = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(10), Inches(0.8))
    title_shape.fill.solid()
    title_shape.fill.fore_color.rgb = PRIMARY_COLOR
    title_shape.line.color.rgb = PRIMARY_COLOR
    
    # Add title text
    title_frame = title_shape.text_frame
    title_frame.clear()
    p = title_frame.paragraphs[0]
    p.text = title
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = WHITE
    title_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    
    return slide

def add_bullet_points(slide, bullets, left=Inches(0.8), top=Inches(1.2), 
                      width=Inches(8.4), height=Inches(5)):
    """Add bullet points to a slide"""
    textbox = slide.shapes.add_textbox(left, top, width, height)
    text_frame = textbox.text_frame
    text_frame.word_wrap = True
    
    for i, bullet in enumerate(bullets):
        if i == 0:
            p = text_frame.paragraphs[0]
        else:
            p = text_frame.add_paragraph()
        
        p.text = bullet
        p.level = 0
        p.font.size = Pt(24)
        p.font.color.rgb = TEXT_COLOR
        p.space_before = Pt(12)
        p.space_after = Pt(12)
    
    return textbox

def add_two_column_slide(prs, title, left_content, right_content, left_title="", right_title=""):
    """Add a two-column content slide"""
    slide = add_content_slide(prs, title)
    
    # Left column
    left_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(4.5), Inches(5.5))
    left_frame = left_box.text_frame
    left_frame.word_wrap = True
    
    if left_title:
        p = left_frame.paragraphs[0]
        p.text = left_title
        p.font.size = Pt(20)
        p.font.bold = True
        p.font.color.rgb = PRIMARY_COLOR
        p.space_after = Pt(12)
    
    for item in left_content:
        p = left_frame.add_paragraph()
        p.text = item
        p.font.size = Pt(18)
        p.font.color.rgb = TEXT_COLOR
        p.space_before = Pt(6)
        p.space_after = Pt(6)
    
    # Right column
    right_box = slide.shapes.add_textbox(Inches(5.2), Inches(1.2), Inches(4.3), Inches(5.5))
    right_frame = right_box.text_frame
    right_frame.word_wrap = True
    
    if right_title:
        p = right_frame.paragraphs[0]
        p.text = right_title
        p.font.size = Pt(20)
        p.font.bold = True
        p.font.color.rgb = PRIMARY_COLOR
        p.space_after = Pt(12)
    
    for item in right_content:
        p = right_frame.add_paragraph()
        p.text = item
        p.font.size = Pt(18)
        p.font.color.rgb = TEXT_COLOR
        p.space_before = Pt(6)
        p.space_after = Pt(6)
    
    return slide

def create_presentation():
    """Create the complete ZipMart investor pitch presentation"""
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    
    # Slide 1: Title Slide
    add_title_slide(prs, "ZipMart", 
                   "Quick Commerce for Tier 3 India\nBlinkit-style delivery in 35 minutes")
    
    # Slide 2: The Problem
    slide = add_content_slide(prs, "The Problem: Underserved Markets")
    bullets = [
        "🏘️ Tier 3 towns (20k-1 lakh population) have no quick commerce options",
        "📱 Poor internet connectivity (2G/3G) limits app usage in these areas",
        "🏪 Existing kirana stores struggle with inventory management & discovery",
        "💰 Majority of consumers prefer Cash on Delivery (low digital payment penetration)",
        "🚚 No efficient last-mile delivery networks in small towns"
    ]
    add_bullet_points(slide, bullets)
    
    # Slide 3: Market Opportunity
    slide = add_content_slide(prs, "Market Opportunity: Massive & Growing")
    bullets = [
        "📊 India has 6,000+ Tier 3 towns with 10+ million small kirana stores",
        "💵 Tier 3 quick commerce market: ₹2,000+ Cr annually (2024-2026 estimates)",
        "📈 Quick commerce segment growing at 45% YoY across India",
        "🌐 Only 5-10% penetration in Tier 3 towns vs. 35%+ in metros",
        "👥 160+ million small-town residents with growing e-commerce adoption"
    ]
    add_bullet_points(slide, bullets)
    
    # Slide 4: The Solution - ZipMart
    slide = add_content_slide(prs, "The Solution: ZipMart")
    bullets = [
        "🤝 Partner with existing local kirana stores (not build dark stores)",
        "⚡ 25-35 minute delivery promise with lean operations",
        "📞 Voice search in Hindi + low-bandwidth optimized app (works on 2G)",
        "💳 Cash on Delivery as primary payment + UPI/Cards as alternatives",
        "🎯 Android-first, optimized for ₹6,000 budget phones (2GB RAM)"
    ]
    add_bullet_points(slide, bullets)
    
    # Slide 5: Business Model
    slide = add_two_column_slide(prs, "Business Model: Multi-Sided Marketplace",
        left_content=[
            "📍 Customers: 3-5% commission per order",
            "🏪 Kirana Stores: 2% platform fee + inventory management tools",
            "🛵 Delivery Partners: 40-50% of delivery fee per order",
            "🎟️ Promotions & advertising: Premium placement fees"
        ],
        right_content=[
            "💰 Revenue streams are sustainable & scalable",
            "✅ Path to profitability in 18 months",
            "📊 Target 20-25% operating margins by Year 3",
            "🚀 Unit economics: ₹40-50 CAC, ₹800+ LTV"
        ],
        left_title="Revenue Sources",
        right_title="Financial Targets"
    )
    
    # Slide 6: Technology Stack
    slide = add_two_column_slide(prs, "Technology Stack: Production-Ready",
        left_content=[
            "🖥️ Backend: Node.js 20 + Express + PostgreSQL 15",
            "📱 Mobile: React Native 0.73 + Expo",
            "🗺️ Maps: Google Maps API + PostGIS for geo-queries",
            "🔐 Auth: Firebase OTP + JWT"
        ],
        right_content=[
            "💳 Payments: Razorpay integration",
            "🔄 Real-time: Socket.io for live tracking",
            "☁️ Cloud: AWS S3 + Redis caching",
            "📢 Notifications: FCM + SMS integration"
        ],
        left_title="Core Technologies",
        right_title="Integrations"
    )
    
    # Slide 7: Key Features
    slide = add_content_slide(prs, "Key Features: Built for Tier 3 India")
    bullets = [
        "🎤 Voice Search in Hindi with auto-transliteration support",
        "📵 Offline mode for critical features (works on 2G/3G)",
        "🗺️ Live delivery tracking with delivery partner location updates",
        "🛍️ Smart inventory sync with kirana stores to prevent overselling",
        "💬 Multi-language support (Hindi, English, regional languages)",
        "📱 Lightweight app (< 20MB after optimization for low-end devices)"
    ]
    add_bullet_points(slide, bullets)
    
    # Slide 8: Go-to-Market Strategy
    slide = add_two_column_slide(prs, "Go-to-Market Strategy",
        left_content=[
            "🎯 Phase 1 (Months 1-3): Single tier-3 city beta (10,000 users)",
            "📍 Focus on kirana partnerships & local marketing",
            "🧪 Validate unit economics & product-market fit",
            "🔄 Rapid iteration based on feedback"
        ],
        right_content=[
            "📈 Phase 2 (Months 4-12): Expand to 10-15 tier-3 cities",
            "🤝 Build supply partnerships with local distributors",
            "💰 Sustainable unit economics at scale",
            "🚀 Phase 3: Build venture-scale platform"
        ],
        left_title="Phase 1: Beta Launch",
        right_title="Phases 2-3: Scale"
    )
    
    # Slide 9: Competitive Advantage
    slide = add_content_slide(prs, "Competitive Advantages: Defensible Moat")
    bullets = [
        "🏪 Kirana partnership model is capital-light vs. dark store competitors",
        "📱 Low-bandwidth tech stack is hard to copy & gives rural reach advantage",
        "🎯 Deep local market expertise in Tier 3 towns (founder background)",
        "🔗 Network effects: more stores → faster delivery → more customers",
        "🌐 Early mover advantage in tier 3 (Blinkit/Zepto focused on metros)"
    ]
    add_bullet_points(slide, bullets)
    
    # Slide 10: Metrics & Traction
    slide = add_two_column_slide(prs, "Metrics & Traction (Beta Phase)",
        left_content=[
            "🎯 Target: 10,000 users by Month 3",
            "📊 Daily active users: 2,000-3,000",
            "💵 Average order value: ₹180-250",
            "⏱️ Average delivery time: 28 minutes"
        ],
        right_content=[
            "📈 Order growth rate: 15-20% MoM",
            "🔄 Repeat customer rate: 45%+",
            "🚀 NPS score target: 65+",
            "✅ Unit economics: Positive by Month 6"
        ],
        left_title="User Metrics",
        right_title="Business Metrics"
    )
    
    # Slide 11: Funding Ask
    slide = add_content_slide(prs, "Funding Ask & Use of Funds")
    
    # Create a simple table-like layout
    textbox = slide.shapes.add_textbox(Inches(1), Inches(1.2), Inches(8), Inches(5.5))
    tf = textbox.text_frame
    tf.word_wrap = True
    
    sections = [
        ("Total Seed Round: ₹2 Crores", True),
        ("", False),
        ("💰 Product & Engineering (45%): ₹90 Lakhs", False),
        ("📱 Mobile app, backend, infrastructure, QA", False),
        ("", False),
        ("📊 Operations & Launch (30%): ₹60 Lakhs", False),
        ("Team hiring, kirana partnerships, initial marketing", False),
        ("", False),
        ("📢 Marketing & GTM (20%): ₹40 Lakhs", False),
        ("User acquisition, brand building in first city", False),
        ("", False),
        ("💼 Admin & Buffer (5%): ₹10 Lakhs", False),
        ("Legal, compliance, contingency", False),
    ]
    
    for text, is_header in sections:
        if text == "":
            p = tf.add_paragraph()
            p.text = ""
            p.space_after = Pt(6)
        else:
            p = tf.add_paragraph()
            p.text = text
            p.font.size = Pt(20 if is_header else 18)
            p.font.bold = is_header
            p.font.color.rgb = PRIMARY_COLOR if is_header else TEXT_COLOR
            p.space_before = Pt(8)
            p.space_after = Pt(6)
    
    # Slide 12: Team
    slide = add_content_slide(prs, "The Team: Built for Execution")
    bullets = [
        "👨‍💼 Founder & CEO: 5+ years in quick commerce & logistics",
        "🔧 CTO: Full-stack engineer with 8+ years experience, shipped 3 apps",
        "📊 COO: Operations expert from Tier 2-3 expansion projects",
        "🤝 Advisor: Former kirana store operator (supply chain insights)",
        "📈 Advisor: Angel investor from Blinkit S1 round"
    ]
    add_bullet_points(slide, bullets)
    
    # Slide 13: Risks & Mitigation
    slide = add_two_column_slide(prs, "Risks & Mitigation Strategies",
        left_content=[
            "⚠️ Kirana adoption: Partner through trusted local operators",
            "📱 User adoption: Voice search & Hindi UI removes friction",
            "🚚 Delivery coverage: Start in high-density neighborhoods",
            "💳 Payments: COD + UPI caters to local preferences"
        ],
        right_content=[
            "🔄 Rapid iteration: 2-week sprint cycles",
            "🎯 Local focus: Single city before scaling",
            "📊 Data-driven: Real-time metrics & A/B testing",
            "💪 Capital efficiency: Bootstrap to unit-positive quickly"
        ],
        left_title="Key Risks",
        right_title="Our Mitigation"
    )
    
    # Slide 14: 18-Month Roadmap
    slide = add_content_slide(prs, "18-Month Roadmap")
    bullets = [
        "📅 Months 1-3: Beta in 1 city, validate product-market fit",
        "📅 Months 4-6: Expand to 3 cities, optimize unit economics",
        "📅 Months 7-12: Scale to 8-10 cities, profitability at unit level",
        "📅 Months 13-18: Expansion roadmap for ₹100 Cr revenue run rate"
    ]
    add_bullet_points(slide, bullets)
    
    # Slide 15: Why ZipMart? Why Now?
    slide = add_content_slide(prs, "Why ZipMart? Why Now?")
    bullets = [
        "📱 Internet penetration in Tier 3 towns crossed 60% (2024)",
        "📦 Logistics costs dropping due to improved infrastructure",
        "🏪 Kirana stores seeking digital transformation (post-COVID)",
        "💰 Quick commerce shift from luxury → essential service",
        "🌐 Proven business model (Blinkit, Zeeker, Swiggy Instamart success)"
    ]
    add_bullet_points(slide, bullets)
    
    # Slide 16: Call to Action
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = PRIMARY_COLOR
    
    # Main message
    cta_box = slide.shapes.add_textbox(Inches(0.5), Inches(2), Inches(9), Inches(2))
    cta_frame = cta_box.text_frame
    cta_frame.word_wrap = True
    
    p = cta_frame.paragraphs[0]
    p.text = "Join us in revolutionizing\nquick commerce for Tier 3 India"
    p.font.size = Pt(52)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER
    
    # Subtitle
    sub_box = slide.shapes.add_textbox(Inches(0.5), Inches(4.2), Inches(9), Inches(2))
    sub_frame = sub_box.text_frame
    sub_frame.word_wrap = True
    
    p = sub_frame.paragraphs[0]
    p.text = "Let's build the future of quick commerce together"
    p.font.size = Pt(32)
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER
    
    # Contact info
    contact_box = slide.shapes.add_textbox(Inches(0.5), Inches(5.8), Inches(9), Inches(1))
    contact_frame = contact_box.text_frame
    contact_frame.word_wrap = True
    
    p = contact_frame.paragraphs[0]
    p.text = "📧 hello@zipmart.in | 📱 +91 XXXXX XXXXX | 🌐 www.zipmart.in"
    p.font.size = Pt(20)
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER
    
    # Save presentation
    output_path = "/Users/anishvrawal/Desktop/retail-sales-data-lake/ZipMart_Investor_Pitch.pptx"
    prs.save(output_path)
    print(f"✅ Presentation created successfully!")
    print(f"📍 Saved to: {output_path}")
    print(f"📊 Total slides: {len(prs.slides)}")

if __name__ == "__main__":
    create_presentation()
